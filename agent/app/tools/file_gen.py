"""File generation tool: materializes outputs as downloadable files."""
import csv
import hashlib
import logging
import os
import re
from datetime import datetime
from pathlib import Path
from typing import Optional

from app.core.config import settings

logger = logging.getLogger(__name__)

FORMAT_EXTENSIONS = {
    "html": "html",
    "htm": "html",
    "md": "md",
    "markdown": "md",
    "csv": "csv",
    "xlsx": "xlsx",
    "xls": "xlsx",
    "doc": "docx",
    "docx": "docx",
    "word": "docx",
    "ppt": "pptx",
    "pptx": "pptx",
}

CONTENT_TYPES = {
    "html": "text/html",
    "md": "text/markdown",
    "csv": "text/csv",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def _detect_format(content: str) -> str:
    """Detect output format from content."""
    if re.search(r'<html|<div|<table|<p\s', content, re.IGNORECASE):
        return "html"
    if re.search(r'^#+\s|\*\*|```', content, re.MULTILINE):
        return "md"
    if re.search(r'^\w+[,\t]', content, re.MULTILINE) and content.count(",") > 5:
        return "csv"
    return "md"


def _normalize_format(fmt: str) -> str:
    return FORMAT_EXTENSIONS.get((fmt or "").strip().lower(), (fmt or "").strip().lower() or "md")


def _generate_filename(session_id: str, fmt: str) -> str:
    """Generate a unique filename."""
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    short_hash = hashlib.md5(session_id.encode()).hexdigest()[:6]
    return f"output_{ts}_{short_hash}.{_normalize_format(fmt)}"


def _strip_code_fence(content: str, fmt: str) -> str:
    pattern = rf"^\s*```(?:{re.escape(fmt)}|markdown|md|html|csv|text)?\s*\n([\s\S]*?)\n```\s*$"
    match = re.match(pattern, content.strip(), re.IGNORECASE)
    return match.group(1) if match else content


def _markdown_rows(content: str) -> list[list[str]]:
    rows: list[list[str]] = []
    for raw in content.splitlines():
        line = raw.strip()
        if not line.startswith("|") or not line.endswith("|"):
            continue
        cells = [cell.strip() for cell in line.strip("|").split("|")]
        if cells and not all(re.fullmatch(r":?-{3,}:?", cell or "") for cell in cells):
            rows.append(cells)
    return rows


def _write_docx(filepath: Path, content: str) -> None:
    from docx import Document

    doc = Document()
    for raw in content.splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue
        heading = re.match(r"^(#{1,6})\s+(.+)$", line)
        if heading:
            doc.add_heading(heading.group(2).strip(), level=min(len(heading.group(1)), 4))
        elif re.match(r"^\s*[-*+]\s+", line):
            doc.add_paragraph(re.sub(r"^\s*[-*+]\s+", "", line).strip(), style="List Bullet")
        elif re.match(r"^\s*\d+[.)]\s+", line):
            doc.add_paragraph(re.sub(r"^\s*\d+[.)]\s+", "", line).strip(), style="List Number")
        else:
            doc.add_paragraph(line)
    doc.save(filepath)


def _write_xlsx(filepath: Path, content: str) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Output"

    rows = _markdown_rows(content)
    if not rows:
        try:
            rows = list(csv.reader(content.splitlines()))
        except csv.Error:
            rows = []
    if not rows:
        rows = [[line] for line in content.splitlines() if line.strip()]

    for row in rows:
        ws.append(row)
    wb.save(filepath)


def _split_slides(content: str) -> list[tuple[str, list[str]]]:
    slides: list[tuple[str, list[str]]] = []
    title = "Output"
    bullets: list[str] = []
    for raw in content.splitlines():
        line = raw.strip()
        if not line:
            continue
        heading = re.match(r"^(?:#{1,3}\s+|Slide\s+\d+[:.-]\s*)(.+)$", line, re.IGNORECASE)
        if heading:
            if bullets or title != "Output":
                slides.append((title, bullets[:]))
            title = heading.group(1).strip()
            bullets = []
        else:
            bullets.append(re.sub(r"^\s*(?:[-*+]|\d+[.)])\s+", "", line))
    if bullets or title:
        slides.append((title, bullets))
    return slides[:30] or [("Output", [content[:800]])]


def _write_pptx(filepath: Path, content: str) -> None:
    from pptx import Presentation

    prs = Presentation()
    for title, bullets in _split_slides(content):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title[:120]
        body = slide.placeholders[1].text_frame
        body.clear()
        for index, bullet in enumerate(bullets[:8]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = bullet[:240]
            paragraph.level = 0
    prs.save(filepath)


async def generate_file(
    content: str,
    session_id: str = "",
    user_id: int = 0,
    format_hint: str = "",
) -> Optional[dict]:
    """Generate a downloadable file from content.

    Returns dict with {filename, url, size, format} or None on failure.
    """
    if not content or len(content) < 100:
        return None

    fmt = _normalize_format(format_hint or _detect_format(content))
    filename = _generate_filename(session_id or "anon", fmt)

    # Ensure output directory exists
    output_dir = Path(settings.GENERATED_FILES_DIR)
    output_dir.mkdir(parents=True, exist_ok=True)

    filepath = output_dir / filename

    try:
        content = _strip_code_fence(content, fmt)

        if fmt == "docx":
            _write_docx(filepath, content)
        elif fmt == "xlsx":
            _write_xlsx(filepath, content)
        elif fmt == "pptx":
            _write_pptx(filepath, content)
        elif fmt == "html" and not content.strip().startswith("<!DOCTYPE"):
            wrapped = f"""<!DOCTYPE html>
<html lang="zh">
<head><meta charset="utf-8"><title>Generated Report</title>
<style>body{{font-family:system-ui;max-width:800px;margin:2em auto;padding:0 1em;line-height:1.6}}</style>
</head>
<body>
{content}
</body>
</html>"""
            filepath.write_text(wrapped, encoding="utf-8")
        else:
            filepath.write_text(content, encoding="utf-8")

        size = filepath.stat().st_size
        # URL assumes Go gateway serves /api/v1/files/
        url = f"/api/v1/files/{filename}"

        logger.info("generated file: %s (%d bytes)", filename, size)
        return {
            "file_id": filename,
            "filename": filename,
            "url": url,
            "download_url": url,
            "size": size,
            "format": fmt,
            "content_type": CONTENT_TYPES.get(fmt, "application/octet-stream"),
        }
    except Exception as e:
        logger.error("file generation failed: %s", e)
        return None
